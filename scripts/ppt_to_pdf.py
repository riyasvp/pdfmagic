#!/usr/bin/env python3
"""
Convert PowerPoint to PDF.
Usage: python ppt_to_pdf.py <input_pptx>
Output: JSON with result
"""

import sys
import os
import json
from datetime import datetime

try:
    from pptx import Presentation
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image as RLImage
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER
    from reportlab.lib.units import inch
except ImportError as e:
    print(json.dumps({"success": False, "error": f"Missing dependency: {str(e)}"}))
    sys.exit(1)

DOWNLOAD_DIR = "/home/z/my-project/download"

def ppt_to_pdf(input_path):
    """Convert PowerPoint to PDF."""
    if not os.path.exists(input_path):
        return {"success": False, "error": f"File not found: {input_path}"}
    
    try:
        # Load PowerPoint
        prs = Presentation(input_path)
        
        # Generate output filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"ppt_to_pdf_{timestamp}.pdf"
        output_path = os.path.join(DOWNLOAD_DIR, output_filename)
        
        # Ensure download directory exists
        os.makedirs(DOWNLOAD_DIR, exist_ok=True)
        
        # Create PDF document
        pdf_doc = SimpleDocTemplate(output_path, pagesize=landscape(A4))
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Heading1'],
            fontSize=24,
            alignment=TA_CENTER,
            spaceAfter=20
        )
        body_style = ParagraphStyle(
            'SlideBody',
            parent=styles['Normal'],
            fontSize=14,
            leading=20,
            alignment=TA_LEFT
        )
        
        story = []
        
        for i, slide in enumerate(prs.slides):
            # Add slide number
            story.append(Paragraph(f"<b>Slide {i + 1}</b>", title_style))
            story.append(Spacer(1, 12))
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            story.append(Paragraph(text, body_style))
                            story.append(Spacer(1, 8))
            
            story.append(PageBreak())
        
        # Remove last page break
        if story and story[-1] == PageBreak():
            story.pop()
        
        pdf_doc.build(story)
        
        return {"success": True, "output": output_path, "slides_converted": len(prs.slides)}
    
    except Exception as e:
        return {"success": False, "error": str(e)}

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "error": "Input PowerPoint file required"}))
        sys.exit(1)
    
    input_path = sys.argv[1]
    result = ppt_to_pdf(input_path)
    print(json.dumps(result))
