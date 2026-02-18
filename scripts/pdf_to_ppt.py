#!/usr/bin/env python3
"""
Convert PDF to PowerPoint.
Usage: python pdf_to_ppt.py <input_pdf>
Output: JSON with result
"""

import sys
import os
import json
from datetime import datetime

try:
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from PIL import Image
except ImportError as e:
    print(json.dumps({"success": False, "error": f"Missing dependency: {str(e)}"}))
    sys.exit(1)

DOWNLOAD_DIR = "/home/z/my-project/download"

def pdf_to_ppt(input_path):
    """Convert PDF pages to PowerPoint slides."""
    if not os.path.exists(input_path):
        return {"success": False, "error": f"File not found: {input_path}"}
    
    try:
        # Convert PDF to images
        images = convert_from_path(input_path, dpi=150)
        
        if not images:
            return {"success": False, "error": "No pages found in PDF"}
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Set slide dimensions to 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Blank layout
        blank_layout = prs.slide_layouts[6]
        
        # Generate output filename for temporary images
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        temp_dir = os.path.join(DOWNLOAD_DIR, f"temp_ppt_{timestamp}")
        os.makedirs(temp_dir, exist_ok=True)
        
        for i, image in enumerate(images):
            # Save image temporarily
            image_path = os.path.join(temp_dir, f"slide_{i+1}.png")
            image.save(image_path, "PNG")
            
            # Add slide
            slide = prs.slides.add_slide(blank_layout)
            
            # Calculate dimensions to fit slide
            img_width, img_height = image.size
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Scale to fit while maintaining aspect ratio
            ratio = min(slide_width / img_width, slide_height / img_height)
            final_width = int(img_width * ratio)
            final_height = int(img_height * ratio)
            
            # Center on slide
            left = (slide_width - final_width) // 2
            top = (slide_height - final_height) // 2
            
            # Add image to slide
            slide.shapes.add_picture(image_path, left, top, final_width, final_height)
        
        # Generate output filename
        output_filename = f"pdf_to_ppt_{timestamp}.pptx"
        output_path = os.path.join(DOWNLOAD_DIR, output_filename)
        
        # Ensure download directory exists
        os.makedirs(DOWNLOAD_DIR, exist_ok=True)
        
        # Save presentation
        prs.save(output_path)
        
        # Clean up temporary files
        import shutil
        shutil.rmtree(temp_dir, ignore_errors=True)
        
        return {
            "success": True,
            "output": output_path,
            "slides_created": len(images)
        }
    
    except Exception as e:
        return {"success": False, "error": str(e)}

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "error": "Input PDF file required"}))
        sys.exit(1)
    
    input_path = sys.argv[1]
    result = pdf_to_ppt(input_path)
    print(json.dumps(result))
